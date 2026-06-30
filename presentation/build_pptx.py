#!/usr/bin/env python3
"""Génère Presentation_PFE_ONMM.pptx — soutenance PFE ONMM.

Implémente la charte et le récit définis dans 00_Conception_Soutenance_ONMM.md.
Chaque slide embarque ses notes orateur (discours + reveal). 16:9.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Charte ----
BLEU   = RGBColor(0x0A, 0x3D, 0x62)
CYAN   = RGBColor(0x1B, 0x9A, 0xAA)
VERT   = RGBColor(0x2E, 0xC4, 0xB6)
CORAIL = RGBColor(0xE6, 0x39, 0x46)
OR     = RGBColor(0xE8, 0xB0, 0x4B)
BLANC  = RGBColor(0xFF, 0xFF, 0xFF)
GRISBG = RGBColor(0xF5, 0xF7, 0xFA)
ANTHRA = RGBColor(0x2D, 0x37, 0x48)
GRIS   = RGBColor(0x71, 0x80, 0x96)

TITLE_FONT = "Montserrat"
BODY_FONT  = "Inter"
MONO_FONT  = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        shadow=False, radius=None):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        # Réutilise l'unique <a:effectLst> créé par shadow.inherit=False
        # (en créer un second corromprait le fichier — 2 effectLst = XML invalide).
        spPr = sp._element.spPr
        eff = spPr.find(qn('a:effectLst'))
        if eff is None:
            eff = spPr.makeelement(qn('a:effectLst'), {})
            spPr.append(eff)
        sh = eff.makeelement(qn('a:outerShdw'),
                             {'blurRad': '90000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = eff.makeelement(qn('a:srgbClr'), {'val': '0A3D62'})
        alpha = eff.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha); sh.append(clr); eff.append(sh)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        font=BODY_FONT, wrap=True):
    """runs: list of (text, size, color, bold, italic) or list of such lists per paragraph."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(para, tuple):
            para = [para]
        for (text, size, color, bold, italic) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = font
    return tb


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def footer(s, n):
    txt(s, 0.4, 7.05, 9, 0.35,
        [("ONMM · PFE 2026 · M. Sidi Mohamed Mohamed Salem", 9, GRIS, False, False)],
        font=BODY_FONT)
    txt(s, 12.3, 7.05, 0.8, 0.35, [(str(n), 9, GRIS, True, False)], align=PP_ALIGN.RIGHT)


def header(s, kicker, title, accent=CYAN):
    box(s, 0.0, 0.0, 0.18, 7.5, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 0.6, 0.45, 11.5, 0.4, [(kicker.upper(), 13, accent, True, False)], font=TITLE_FONT)
    txt(s, 0.6, 0.85, 12.1, 1.0, [(title, 30, BLEU, True, False)], font=TITLE_FONT)
    box(s, 0.62, 1.72, 1.6, 0.05, fill=accent, shape=MSO_SHAPE.RECTANGLE)


def chip(s, l, t, w, h, label, color, fg=BLANC, size=12.5, bold=True):
    box(s, l, t, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, l, t, w, h, [(label, size, fg, bold, False)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)


def card(s, l, t, w, h, title, body, accent=CYAN, title_size=16, body_size=12.5):
    box(s, l, t, w, h, fill=BLANC, line=None, shadow=True, radius=0.07)
    box(s, l, t, w, 0.12, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    txt(s, l + 0.25, t + 0.28, w - 0.5, 0.6, [(title, title_size, BLEU, True, False)], font=TITLE_FONT)
    if body:
        txt(s, l + 0.25, t + 0.95, w - 0.5, h - 1.1,
            [[(b, body_size, ANTHRA, False, False)] for b in body])


# ============================================================ SLIDE 1 — Couverture
s = slide(); bg(s, BLEU)
box(s, 0, 6.6, 13.333, 0.9, fill=CYAN, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 6.55, 13.333, 0.06, fill=VERT, shape=MSO_SHAPE.RECTANGLE)
txt(s, 0.9, 0.7, 11.5, 0.5, [("PROJET DE FIN D'ÉTUDES · LICENCE DAII · 2025–2026", 14, CYAN, True, False)], font=TITLE_FONT)
txt(s, 0.9, 1.7, 11.6, 2.4,
    [("Conception & développement d'une application web", 38, BLANC, True, False)],
    font=TITLE_FONT)
txt(s, 0.9, 2.75, 11.6, 1.0, [("pour l'Ordre National des Médecins de Mauritanie", 30, VERT, True, False)], font=TITLE_FONT)
txt(s, 0.95, 4.0, 11.0, 0.6,
    [("Une plateforme centralisée, sécurisée et traçable au service de la régulation médicale", 17, RGBColor(0xC9,0xE6,0xEC), False, True)])
txt(s, 0.95, 5.15, 7, 1.2,
    [[("Présenté par : ", 14, BLANC, True, False), ("Mohamed Sidi Mohamed Mohamed Salem  (C23282)", 14, BLANC, False, False)],
     [("Encadrants : ", 13, RGBColor(0xC9,0xE6,0xEC), True, False), ("Dr. Ahmed Sejad · Ing. Mohamed Saleck Elbechir", 13, RGBColor(0xC9,0xE6,0xEC), False, False)],
     [("Entreprise d'accueil : ", 13, RGBColor(0xC9,0xE6,0xEC), True, False), ("SMART MS SA", 13, RGBColor(0xC9,0xE6,0xEC), False, False)]])
txt(s, 9.0, 6.7, 4.0, 0.7, [("Soutenance : 11/07/2026", 14, BLEU, True, False)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
notes(s, "ACCUEIL (35s). Se présenter, annoncer le titre. Ton : projet mené EN ENTREPRISE sur un besoin RÉEL. Ne pas lire les noms d'encadrants. Enchaîner : « voici en une minute le chemin que je vais vous faire suivre ». Aucune animation.")

# ============================================================ SLIDE 2 — Fil conducteur
s = slide(); bg(s, BLANC)
header(s, "Fil conducteur", "Le fil de cette présentation")
steps = [("1", "Contexte\n& mission", CYAN), ("2", "Le\nproblème", CORAIL), ("3", "La\nsolution", VERT),
         ("4", "La\nconstruction", BLEU), ("5", "La\nsécurité", OR), ("6", "Bilan &\nperspectives", CYAN)]
x = 0.7; w = 1.95; gap = 0.12
for i, (num, label, col) in enumerate(steps):
    lx = x + i * (w + gap)
    box(s, lx, 3.2, w, 1.7, fill=GRISBG, shadow=True, radius=0.08)
    chip(s, lx + w/2 - 0.35, 3.45, 0.7, 0.7, num, col, size=20)
    txt(s, lx, 4.25, w, 0.7, [[(p, 12.5, BLEU, True, False)] for p in label.split("\n")],
        align=PP_ALIGN.CENTER, font=TITLE_FONT)
    if i < 5:
        box(s, lx + w + gap/2 - 0.06, 3.95, 0.18, 0.18, fill=GRIS, shape=MSO_SHAPE.CHEVRON)
txt(s, 0.7, 5.5, 12, 0.6,
    [[("Un récit, pas un sommaire : ", 15, ANTHRA, True, False),
      ("un problème réel → une solution argumentée → un climax technique (la sécurité du vote).", 15, ANTHRA, False, False)]])
footer(s, 2)
notes(s, "FIL (40s). Annoncer le VOYAGE, pas les chapitres. Insister : le jalon 5 (sécurité) est le cœur technique. Reveal : les 6 jalons un par un. Ne pas détailler. Transition → « commençons par pourquoi ce projet existe ».")

# ============================================================ SLIDE 3 — Contexte
s = slide(); bg(s, BLANC)
header(s, "1 · Contexte", "Un projet né d'une priorité nationale")
card(s, 0.7, 2.1, 3.9, 2.6, "Le cadre national",
     ["Stratégie nationale", "e-Santé 2024–2030 :", "numériser les services", "médicaux & administratifs."], accent=CYAN)
card(s, 4.75, 2.1, 3.9, 2.6, "L'entreprise",
     ["SMART MS SA — éditeur", "mauritanien, e-santé &", "digitalisation publique", "(Odoo Gold Partner).", "Stage de 4 mois."], accent=BLEU)
card(s, 8.8, 2.1, 3.85, 2.6, "Le client : l'ONMM",
     ["Autorité ordinale créée", "en 2019 (décret 2019-077).", "≈ 1 028 médecins inscrits.", "Inscription OBLIGATOIRE", "pour exercer légalement."], accent=VERT)
# frise
box(s, 0.7, 5.35, 11.95, 0.9, fill=GRISBG, radius=0.2)
fr = [("2019", "Création ONMM"), ("Janv. 2024", "1ᵉʳ congrès national"), ("2024–2030", "Stratégie e-Santé"), ("2026", "Notre projet")]
for i, (d, l) in enumerate(fr):
    lx = 1.0 + i * 2.95
    chip(s, lx, 5.5, 1.2, 0.4, d, BLEU, size=12)
    txt(s, lx + 1.3, 5.5, 1.6, 0.4, [(l, 10.5, ANTHRA, False, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)
notes(s, "CONTEXTE (1min). Raconter la dynamique e-santé. CHIFFRE CLÉ à dire : 1028 médecins + inscription obligatoire = enjeu de masse ET légal. Ne pas lire la déf. d'Odoo Partner. Q probable : « pourquoi l'ONMM pas le ministère ? » → l'ONMM est l'autorité ordinale, détenteur légitime du tableau des médecins.")

# ============================================================ SLIDE 4 — ONMM aujourd'hui
s = slide(); bg(s, BLANC)
header(s, "2 · Le problème", "Aujourd'hui : papier, Excel et WhatsApp", accent=CORAIL)
items = [("Adhésion", "Dossier PAPIER au siège,\ndiplômes vérifiés en présentiel"),
         ("Registre", "Fichier EXCEL\nnon protégé"),
         ("Suivi dossier", "Téléphone / e-mail / WhatsApp\nou déplacement à Nouakchott"),
         ("Réclamations", "Archives PAPIER\nnon exploitables"),
         ("Élections", "100 % MANUELLES")]
for i, (t, d) in enumerate(items):
    col = i % 3; row = i // 3
    lx = 0.7 + col * 4.05; ly = 2.2 + row * 1.55
    box(s, lx, ly, 3.85, 1.35, fill=RGBColor(0xFD,0xEC,0xEA), line=CORAIL, line_w=1.0, radius=0.08)
    txt(s, lx + 0.25, ly + 0.15, 3.4, 0.4, [(t, 15, CORAIL, True, False)], font=TITLE_FONT)
    txt(s, lx + 0.25, ly + 0.62, 3.5, 0.7, [[(p, 11.5, ANTHRA, False, False)] for p in d.split("\n")])
box(s, 8.8, 3.75, 3.85, 1.35, fill=BLEU, radius=0.08)
txt(s, 9.0, 3.95, 3.5, 1.0,
    [[("Tout converge", 16, BLANC, True, False)], [("physiquement vers Nouakchott", 14, VERT, True, False)]],
    anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
footer(s, 4)
notes(s, "EXISTANT (1min). Rendre le problème TANGIBLE. Image forte : « un médecin de Néma doit appeler, écrire sur WhatsApp ou faire 1200 km pour connaître l'état de son dossier ». Reveal : canaux un à un, puis surligner « tout converge vers Nouakchott ». Ne pas mépriser les agents — on constate.")

# ============================================================ SLIDE 5 — Diagnostic
s = slide(); bg(s, BLANC)
header(s, "2 · Le problème", "Le vrai problème : l'absence de système d'information", accent=CORAIL)
rows = [("Procédures manuelles", "Retards, erreurs de saisie, temps agent gaspillé"),
        ("Suivi dispersé", "Aucun historique, aucune visibilité (médecin ET admin)"),
        ("Données Excel", "Aucune traçabilité, aucun contrôle d'accès → risque légal"),
        ("Réclamations papier", "Impossibles à exploiter collectivement"),
        ("Aucun espace médecin", "Tout repose sur l'informel"),
        ("Élections manuelles", "Participation des régions éloignées limitée")]
ty = 2.05
txt(s, 0.7, ty, 4.6, 0.4, [("CAUSE", 13, CORAIL, True, False)], font=TITLE_FONT)
txt(s, 6.0, ty, 6.6, 0.4, [("CONSÉQUENCE", 13, ANTHRA, True, False)], font=TITLE_FONT)
for i, (c, cons) in enumerate(rows):
    ly = 2.5 + i * 0.72
    hot = (i == 2)
    box(s, 0.7, ly, 4.9, 0.62, fill=(CORAIL if hot else RGBColor(0xFD,0xEC,0xEA)), radius=0.15)
    txt(s, 0.85, ly, 4.6, 0.62, [(c, 13, (BLANC if hot else CORAIL), True, False)], anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
    box(s, 5.7, ly + 0.21, 0.3, 0.2, fill=GRIS, shape=MSO_SHAPE.CHEVRON)
    box(s, 6.0, ly, 6.65, 0.62, fill=(RGBColor(0xFB,0xE3,0xE0) if hot else GRISBG),
        line=(CORAIL if hot else None), line_w=1.2, radius=0.15)
    txt(s, 6.2, ly, 6.3, 0.62, [(cons, 12.5, ANTHRA, hot, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)
notes(s, "DIAGNOSTIC (1min15) — slide la plus dense, REVEAL LIGNE PAR LIGNE OBLIGATOIRE. Ne pas lire les 6 lignes : en développer 2 (traçabilité + régions éloignées), survoler le reste. Pivot : « pas des dysfonctionnements isolés, mais le symptôme d'une organisation sans outil numérique ». La ligne 3 (rouge) est le pont vers la sécurité. Q : « avez-vous quantifié ? » → constats qualitatifs, limite de l'étude.")

# ============================================================ SLIDE 6 — Problématique
s = slide(); bg(s, BLEU)
box(s, 0.0, 0.0, 0.18, 7.5, fill=CYAN, shape=MSO_SHAPE.RECTANGLE)
txt(s, 0.7, 0.55, 11, 0.5, [("LA QUESTION CENTRALE", 15, CYAN, True, False)], font=TITLE_FONT)
box(s, 1.1, 1.6, 11.1, 2.5, fill=RGBColor(0x0E,0x4F,0x7C), shadow=True, radius=0.06)
box(s, 1.1, 1.6, 0.14, 2.5, fill=VERT, shape=MSO_SHAPE.RECTANGLE)
txt(s, 1.5, 1.9, 10.4, 1.95,
    [("Comment doter l'ONMM d'une plateforme centralisée, sécurisée et accessible "
      "qui structure ses démarches, assure le suivi et la traçabilité des dossiers, "
      "et fluidifie les échanges entre l'administration, les médecins et le public ?",
      22, BLANC, False, False)], anchor=MSO_ANCHOR.MIDDLE)
objs = [("Centraliser", CYAN), ("Dématérialiser\n& suivre", VERT), ("Sécuriser\n& tracer", OR), ("Rapprocher\nrégions / public", CYAN)]
for i, (o, c) in enumerate(objs):
    lx = 0.95 + i * 3.0
    box(s, lx, 4.7, 2.75, 1.6, fill=BLANC, shadow=True, radius=0.1)
    box(s, lx, 4.7, 2.75, 0.12, fill=c, shape=MSO_SHAPE.RECTANGLE)
    chip(s, lx + 0.2, 4.95, 0.55, 0.55, str(i+1), c, size=17)
    txt(s, lx + 0.85, 4.95, 1.8, 1.1, [[(p, 14, BLEU, True, False)] for p in o.split("\n")],
        anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
txt(s, 0.7, 6.55, 11, 0.4, [("Ces 4 objectifs serviront de grille d'évaluation en conclusion.", 13, RGBColor(0xC9,0xE6,0xEC), False, True)])
notes(s, "PROBLÉMATIQUE (50s). LIRE lentement la problématique (seule chose qu'on lit volontairement = c'est le contrat). Annoncer les 4 objectifs comme grille d'évaluation finale. Reveal : problématique seule d'abord, puis les 4 piliers. Ne pas multiplier les objectifs.")

# ============================================================ SLIDE 7 — Benchmark
s = slide(); bg(s, BLANC)
header(s, "2 · Le problème", "Ce que font les autres ordres — et ce que nous ajoutons")
for i, (pays, desc) in enumerate([("France", "Espace membre sécurisé\n+ tableau de l'Ordre"),
                                   ("Maroc", "Info institutionnelle\n+ tableau (contexte proche)"),
                                   ("Sénégal", "Communication\ninstitutionnelle")]):
    lx = 0.7 + i * 4.05
    card(s, lx, 2.05, 3.85, 1.95, pays, desc.split("\n"), accent=BLEU, title_size=17, body_size=12.5)
box(s, 0.7, 4.4, 11.95, 1.95, fill=GRISBG, radius=0.06)
txt(s, 0.95, 4.55, 11, 0.4, [("Échelle de maturité numérique", 13, GRIS, True, False)], font=TITLE_FONT)
seg = [("Portail public", BLEU), ("Annuaire", BLEU), ("Espace membre", BLEU),
       ("Réclamations · Sondages · Élections en ligne", VERT)]
lx = 0.95
for i, (lab, c) in enumerate(seg):
    w = 2.0 if i < 3 else 4.6
    box(s, lx, 5.15, w, 0.85, fill=c, radius=0.12)
    txt(s, lx, 5.15, w, 0.85, [(lab, 12 if i<3 else 12.5, BLANC, True, False)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
    if i < 3:
        box(s, lx + w - 0.02, 5.45, 0.22, 0.22, fill=GRIS, shape=MSO_SHAPE.CHEVRON)
    lx += w + 0.18
txt(s, 0.95, 6.05, 11, 0.3, [("← Repris des bonnes pratiques        Notre valeur ajoutée ONMM →", 11, GRIS, False, True)])
footer(s, 7)
notes(s, "BENCHMARK (50s). Montrer une démarche d'étude comparative. INSISTER : réclamations + sondages + élections n'ont PAS été trouvés sur ces plateformes → valeur ajoutée ONMM (segment vert). Reveal : 3 pays → jauge → segment vert s'illumine. Q : « testé en profondeur ? » → consultation des fonctions publiques, benchmark de cadrage.")

# ============================================================ SLIDE 8 — Solution 360
s = slide(); bg(s, BLANC)
header(s, "3 · La solution", "Une plateforme, trois espaces, huit modules", accent=VERT)
# hub central
box(s, 5.35, 3.3, 2.65, 1.3, fill=BLEU, shadow=True, radius=0.12)
txt(s, 5.35, 3.3, 2.65, 1.3, [[("Plateforme", 16, BLANC, True, False)], [("ONMM", 16, VERT, True, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
# 3 espaces
for i, (esp, c) in enumerate([("Espace PUBLIC", CYAN), ("Espace MÉDECIN", VERT), ("Espace ADMIN", BLEU)]):
    lx = 0.9 + i * 4.0
    chip(s, lx, 2.15, 3.55, 0.6, esp, c, size=14)
# modules en 2 rangées
mods = [("Adhésions", VERT), ("Registre / Annuaire", CYAN), ("Réclamations", CYAN), ("Annonces", CYAN),
        ("Sondages", VERT), ("Élections & vote", CORAIL), ("Notifications", CYAN), ("Spécialités", BLEU)]
for i, (m, c) in enumerate(mods):
    col = i % 4; row = i // 4
    lx = 0.9 + col * 3.05; ly = 5.0 + row * 0.95
    highlight = (m == "Élections & vote")
    box(s, lx, ly, 2.85, 0.78, fill=(c if highlight else BLANC), line=(None if highlight else c),
        line_w=1.4, shadow=highlight, radius=0.18)
    txt(s, lx, ly, 2.85, 0.78, [(m, 12.5, (BLANC if highlight else c), True, False)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
footer(s, 8)
notes(s, "SOLUTION 360 (1min). Vue d'ensemble en un coup d'œil. 3 espaces = 3 portes d'une même maison. TEASER : « le module Élections concentre l'essentiel de la difficulté technique, j'y reviens » (halo sur Élections). Reveal : centre → 3 espaces → modules. Ne pas détailler chaque module (on les voit en images).")

# ============================================================ SLIDE 9 — Acteurs & besoins
s = slide(); bg(s, BLANC)
header(s, "3 · La solution", "Trois acteurs, des exigences claires")
personas = [("Visiteur public", CYAN, ["Consulter annonces", "Rechercher / vérifier", "un médecin", "Déposer une réclamation"]),
            ("Médecin", VERT, ["Adhérer en ligne", "Suivre ses dossiers", "Répondre aux sondages", "Voter aux élections"]),
            ("Administrateur", BLEU, ["Valider les adhésions", "Gérer le registre", "Traiter réclamations", "Organiser élections"])]
for i, (name, c, acts) in enumerate(personas):
    lx = 0.7 + i * 4.05
    box(s, lx, 2.05, 3.85, 2.85, fill=BLANC, shadow=True, radius=0.06)
    chip(s, lx, 2.05, 3.85, 0.6, name, c, size=15)
    txt(s, lx + 0.3, 2.85, 3.3, 1.9, [[("•  " + a, 12.5, ANTHRA, False, False)] for a in acts])
box(s, 0.7, 5.2, 11.95, 1.45, fill=BLEU, radius=0.06)
txt(s, 0.95, 5.32, 11, 0.4, [("Exigences non fonctionnelles", 13, VERT, True, False)], font=TITLE_FONT)
nfr = ["Sécurité", "Traçabilité", "Performance", "Disponibilité", "Ergonomie", "Compatibilité", "Évolutivité"]
lx = 0.95
for i, n in enumerate(nfr):
    em = n in ("Sécurité", "Traçabilité")
    w = 1.55 if not em else 1.7
    box(s, lx, 5.85, w, 0.6, fill=(VERT if em else RGBColor(0x0E,0x4F,0x7C)), radius=0.25)
    txt(s, lx, 5.85, w, 0.6, [(n, 11.5, BLANC, em, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lx += w + 0.12
footer(s, 9)
notes(s, "ACTEURS (45s). Lier chaque acteur à 2 actions max. Souligner : sécurité & traçabilité ne sont pas des options mais des exigences (surlignées en vert). Q : « évolutivité ? » → couches + REST stateless + modules découplés, ajout d'un module sans toucher au noyau.")

# ============================================================ SLIDE 10 — Méthodologie
s = slide(); bg(s, BLANC)
header(s, "4 · La construction", "UML pour concevoir, Scrum pour livrer")
card(s, 0.7, 2.05, 4.0, 3.0, "Modélisation UML",
     ["Cas d'utilisation", "→ fonctions par acteur", "", "Séquence", "→ scénarios clés", "", "Classes", "→ modèle de données"],
     accent=CYAN, body_size=12.5)
txt(s, 5.0, 2.05, 7.6, 0.45, [("Agile Scrum · 8 sprints × 2 semaines · suivi Jira", 14, BLEU, True, False)], font=TITLE_FONT)
sprints = [("S0", "Analyse / UML"), ("S1-2", "Auth · Adhésion · Registre"), ("S3", "Portail · Annuaire"),
           ("S4", "Réclam. · Notifs"), ("S5", "Sondages"), ("S6", "Élections / Vote"), ("S7", "Intégration · Tests")]
for i, (sp, lab) in enumerate(sprints):
    ly = 2.6 + i * 0.55
    c = CORAIL if sp == "S6" else CYAN
    chip(s, 5.0, ly, 0.95, 0.45, sp, c, size=12)
    box(s, 6.05, ly + 0.05, 6.4, 0.35, fill=GRISBG, radius=0.3)
    txt(s, 6.2, ly, 6.2, 0.45, [(lab, 12, ANTHRA, (sp == "S6"), False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)
notes(s, "MÉTHODO (1min). ASSUMER le Scrum solo D'AVANCE : « projet individuel, j'ai gardé les principes (backlog, sprints, incréments, Jira) et adapté les cérémonies » → désamorce la question piège. Reveal : UML puis frise sprint par sprint (matérialise l'itération). S6 (vote) en rouge = annonce le climax.")

# ============================================================ SLIDE 11 — Architecture
s = slide(); bg(s, BLANC)
header(s, "4 · La construction", "Une architecture en couches, découplée et sécurisée")
layers = [("PRÉSENTATION", CYAN, "React.js  ·  React Router  ·  Axios  ·  Tailwind CSS"),
          ("APPLICATION / MÉTIER", BLEU, "Spring Boot (Java)  ·  API REST  ·  Spring Security + JWT (stateless)  ·  RBAC ADMIN / MEDECIN"),
          ("DONNÉES", VERT, "PostgreSQL  ·  relationnel  ·  intégrité référentielle  ·  transactions ACID")]
for i, (name, c, tech) in enumerate(layers):
    ly = 2.15 + i * 1.4
    box(s, 0.7, ly, 11.0, 1.2, fill=BLANC, line=c, line_w=1.5, shadow=True, radius=0.06)
    box(s, 0.7, ly, 0.16, 1.2, fill=c, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 1.05, ly + 0.18, 4.0, 0.8, [(name, 15, c, True, False)], font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 4.6, ly + 0.18, 6.9, 0.85, [(tech, 12.5, ANTHRA, False, False)], anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        box(s, 11.85, ly + 0.45, 0.5, 0.45, fill=c, shape=MSO_SHAPE.DOWN_ARROW)
txt(s, 0.7, 6.55, 11.5, 0.4,
    [[("Flux : ", 12, GRIS, True, False),
      ("Navigateur → HTTPS → API REST → service métier → JPA → PostgreSQL", 12, ANTHRA, False, False)]])
footer(s, 11)
notes(s, "ARCHITECTURE (1min15). NE PAS faire un slide à logos : justifier CHAQUE choix en 1 phrase (cf. tableau justifications). Insister sur STATELESS (= scalabilité, pas de session à voler) et séparation des couches (= maintenabilité/sécurité). Reveal : couche par couche, puis flux animé. Q : « Spring vs Node ? Postgres vs Mongo ? » → réponses préparées.")

# ============================================================ SLIDE 12 — Workflow adhésion
s = slide(); bg(s, BLANC)
header(s, "4 · La construction", "Du dossier papier au compte activé — en ligne")
flow = [("Médecin", "Dépose la demande\n(5 étapes guidées)", VERT),
        ("Admin", "Instruit\n& vérifie", BLEU),
        ("Validation", "Décision\n+ e-mail d'activation", CYAN),
        ("Médecin", "Confirme e-mail +\nmot de passe (BCrypt)", VERT),
        ("Système", "Compte actif +\nregistre officiel", BLEU)]
for i, (who, what, c) in enumerate(flow):
    lx = 0.6 + i * 2.5
    box(s, lx, 2.5, 2.25, 1.7, fill=BLANC, line=c, line_w=1.5, shadow=True, radius=0.08)
    box(s, lx, 2.5, 2.25, 0.45, fill=c, shape=MSO_SHAPE.RECTANGLE)
    txt(s, lx, 2.5, 2.25, 0.45, [(who, 12.5, BLANC, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
    txt(s, lx + 0.12, 3.05, 2.0, 1.1, [[(p, 11.5, ANTHRA, False, False)] for p in what.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        box(s, lx + 2.27, 3.15, 0.22, 0.35, fill=c, shape=MSO_SHAPE.CHEVRON)
box(s, 2.0, 4.7, 5.5, 0.7, fill=RGBColor(0xFD,0xEC,0xEA), line=CORAIL, line_w=1.0, radius=0.2)
txt(s, 2.2, 4.7, 5.2, 0.7, [[("Branche Rejet : ", 12, CORAIL, True, False), ("notification motivée au médecin", 12, ANTHRA, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, 7.9, 4.7, 4.7, 0.7, fill=RGBColor(0xE8,0xF7,0xF5), line=VERT, line_w=1.0, radius=0.2)
txt(s, 8.1, 4.7, 4.3, 0.7, [[("Chaque transition d'état est ", 12, ANTHRA, False, False), ("historisée", 12, VERT, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.7, 5.8, 11.5, 0.7,
    [[("Deux points clés : ", 14, BLEU, True, False),
      ("la traçabilité (chaque changement de statut est tracé) et la sécurité (mot de passe jamais en clair, BCrypt).", 14, ANTHRA, False, False)]])
footer(s, 12)
notes(s, "WORKFLOW (1min). Raconter comme une histoire utilisateur. Pointer TRAÇABILITÉ (chaque transition historisée) et BCrypt (mot de passe jamais en clair) → transition vers sécurité. NE PAS afficher le diagramme de classes complet (annexe). Préférer ce parcours à étapes plutôt qu'un diagramme de séquence brut.")

# ============================================================ SLIDE 13 — Sécurité fondations
s = slide(); bg(s, BLEU)
box(s, 0.0, 0.0, 0.18, 7.5, fill=OR, shape=MSO_SHAPE.RECTANGLE)
txt(s, 0.6, 0.45, 11.5, 0.4, [("5 · LA SÉCURITÉ", 13, OR, True, False)], font=TITLE_FONT)
txt(s, 0.6, 0.85, 12.1, 0.9, [("Sécurité by design : identité, accès, secrets", 30, BLANC, True, False)], font=TITLE_FONT)
box(s, 0.62, 1.72, 1.6, 0.05, fill=OR, shape=MSO_SHAPE.RECTANGLE)
cards = [("Authentification", "JWT", ["Jeton signé renvoyé", "au login", "Politique stateless :", "aucune session serveur", "Validé à chaque requête"], CYAN),
         ("Autorisation", "RBAC · Spring Security", ["Rôles ADMIN / MEDECIN", "Chaque route filtrée", "Principe du moindre", "privilège"], VERT),
         ("Secrets", "BCrypt + Salt", ["Hachage adaptatif", "Sel aléatoire par mdp", "2 mots de passe égaux", "→ empreintes différentes", "Anti rainbow-table"], OR)]
for i, (t, sub, body, c) in enumerate(cards):
    lx = 0.7 + i * 4.05
    box(s, lx, 2.2, 3.85, 4.2, fill=RGBColor(0x0E,0x4F,0x7C), shadow=True, radius=0.06)
    box(s, lx, 2.2, 3.85, 0.12, fill=c, shape=MSO_SHAPE.RECTANGLE)
    txt(s, lx + 0.3, 2.5, 3.3, 0.5, [(t, 17, BLANC, True, False)], font=TITLE_FONT)
    chip(s, lx + 0.3, 3.05, 3.25, 0.55, sub, c, size=13)
    txt(s, lx + 0.35, 3.85, 3.2, 2.4, [[("•  " + b, 12.5, RGBColor(0xD8,0xE8,0xEE), False, False)] for b in body])
notes(s, "SÉCURITÉ FONDATIONS (1min). Rampe de lancement du slide vote. Expliquer STATELESS (pourquoi : scalabilité + pas de session à voler) et SEL (anti rainbow-table). Bref, net, confiant. Reveal : 3 cartes en cascade. Q : « où stockez-vous le JWT ? durée ? refresh ? » → réponse préparée. Transition : « pour le vote, on va beaucoup plus loin ».")

# ============================================================ SLIDE 14 — Vote (CLIMAX)
s = slide(); bg(s, BLEU)
box(s, 0.0, 0.0, 0.18, 7.5, fill=CORAIL, shape=MSO_SHAPE.RECTANGLE)
txt(s, 0.6, 0.4, 11.5, 0.4, [("5 · LA SÉCURITÉ · POINT FORT TECHNIQUE", 13, CORAIL, True, False)], font=TITLE_FONT)
txt(s, 0.6, 0.78, 12.3, 0.8, [("Vote électronique : 4 garanties, 5 mécanismes", 28, BLANC, True, False)], font=TITLE_FONT)
# pipeline
pipe = [("Vote", VERT), ("Chiffrement\nAES-256-GCM", CYAN), ("Clé scellée\nRSA-2048-OAEP", OR),
        ("Hash\nSHA-256", CYAN), ("Signature\nEd25519", VERT)]
for i, (lab, c) in enumerate(pipe):
    lx = 0.6 + i * 2.45
    box(s, lx, 1.95, 2.2, 1.15, fill=RGBColor(0x0E,0x4F,0x7C), line=c, line_w=1.5, radius=0.1)
    txt(s, lx, 1.95, 2.2, 1.15, [[(p, 12, BLANC, True, False)] for p in lab.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
    if i < 4:
        box(s, lx + 2.2, 2.35, 0.25, 0.35, fill=c, shape=MSO_SHAPE.CHEVRON)
txt(s, 0.6, 3.25, 12, 0.35, [("Clé privée RSA verrouillée pendant le scrutin → déverrouillée au dépouillement seulement", 12.5, OR, False, True)])
# tableau propriété/mécanisme
pairs = [("Anonymat de l'électeur", "HMAC-SHA256 → Voter Token (secret+electionId+médecinId)"),
         ("Confidentialité du bulletin", "AES-256-GCM, clé (DEK) aléatoire par bulletin"),
         ("Protection de la clé", "RSA-2048-OAEP : clé privée inaccessible pendant le vote"),
         ("Intégrité des données", "SHA-256 Vote Hash (sur le bulletin déjà chiffré)"),
         ("Authenticité des résultats", "Signature Ed25519 du serveur sur chaque bulletin")]
for i, (prop, mec) in enumerate(pairs):
    ly = 3.75 + i * 0.66
    box(s, 0.6, ly, 4.3, 0.56, fill=CORAIL if i in (0,2) else RGBColor(0x12,0x5A,0x88), radius=0.15)
    txt(s, 0.75, ly, 4.1, 0.56, [(prop, 12, BLANC, True, False)], anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
    box(s, 5.0, ly, 7.6, 0.56, fill=BLANC, radius=0.1)
    txt(s, 5.2, ly, 7.3, 0.56, [(mec, 11.5, ANTHRA, False, False)], anchor=MSO_ANCHOR.MIDDLE, font=MONO_FONT)
notes(s, "VOTE (2min) — SOMMET DRAMATIQUE, ralentir. À DÉVELOPPER : (1) chiffrement HYBRIDE — pourquoi pas tout en RSA : trop lent/limité → AES pour les données, RSA pour la clé. (2) clé privée VERROUILLÉE pendant le vote = personne, pas même l'admin, ne lit un bulletin avant clôture. (3) hash (détecte modif après coup) vs signature (prouve l'origine). NE PAS lire les algos un par un : MONTRER + expliquer la logique. Si on ne sait pas un détail : « le principe est X, l'implémentation est dans le code ». Honnêteté > bluff. Reveal : bulletin traverse le pipeline étape par étape, le verrou s'ouvre au dépouillement.")

# ============================================================ SLIDE 15 — Plateforme en action
s = slide(); bg(s, BLANC)
header(s, "3 · La solution", "La plateforme en action", accent=VERT)
shots = [("Portail public", "Accueil + annuaire des médecins", CYAN),
         ("Adhésion en ligne", "Formulaire multi-étapes guidé", VERT),
         ("Espace médecin", "Tableau de bord personnel", VERT),
         ("Vote & résultats", "Bulletin chiffré + dépouillement", CORAIL),
         ("Espace admin", "Gestion adhésions / élections", BLEU),
         ("Sondages", "Création & statistiques", CYAN)]
for i, (t, d, c) in enumerate(shots):
    col = i % 3; row = i // 3
    lx = 0.7 + col * 4.05; ly = 2.1 + row * 2.35
    box(s, lx, ly, 3.85, 2.1, fill=GRISBG, shadow=True, radius=0.05)
    box(s, lx + 0.15, ly + 0.15, 3.55, 1.25, fill=BLANC, line=c, line_w=1.2, radius=0.04)
    txt(s, lx + 0.15, ly + 0.15, 3.55, 1.25, [("[ capture d'écran ]", 11, GRIS, False, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + 0.25, ly + 1.5, 3.4, 0.35, [(t, 13.5, c, True, False)], font=TITLE_FONT)
    txt(s, lx + 0.25, ly + 1.82, 3.4, 0.3, [(d, 10.5, ANTHRA, False, False)])
footer(s, 15)
notes(s, "RÉALISATION (1min15 ou +90s démo). Remplacer les 6 placeholders par vos VRAIES captures (cf. rapport fig 5.x). Ne pas décrire chaque écran : raconter 2 écrans (adhésion + vote). RECOMMANDATION FORTE : prévoir une DÉMO LIVE courte du vote (effet « ça marche vraiment » imbattable) avec captures en secours. 6 captures max, lisibles.")

# ============================================================ SLIDE 16 — Bilan
s = slide(); bg(s, BLANC)
header(s, "6 · Bilan", "Objectifs atteints — et la suite")
cols = [("RÉALISÉ", VERT, ["Plateforme centralisée", "Adhésions dématérialisées + suivi", "Registre / annuaire public", "Réclamations numériques", "Sondages & statistiques", "Élections sécurisées (crypto)", "RBAC / JWT / BCrypt"]),
        ("LIMITES ASSUMÉES", OR, ["Clés crypto en base (pas de HSM)", "Pas d'audit externe indépendant", "Tests surtout fonctionnels (Postman)", "Déploiement mono-instance"]),
        ("PERSPECTIVES", CYAN, ["HSM + audit externe + archi distribuée", "Tableaux de bord & statistiques", "Notifications avancées", "API officielle du registre national", "(ONMM = source de référence santé)"])]
for i, (t, c, items) in enumerate(cols):
    lx = 0.7 + i * 4.05
    box(s, lx, 2.05, 3.85, 4.6, fill=BLANC, shadow=True, radius=0.05)
    chip(s, lx, 2.05, 3.85, 0.6, t, c, size=14)
    txt(s, lx + 0.28, 2.85, 3.35, 3.6, [[("•  " + it, 12, ANTHRA, False, False)] for it in items])
footer(s, 16)
notes(s, "BILAN (1min). BOUCLE NARRATIVE : revenir aux 4 objectifs du slide 6 (très apprécié). Présenter les limites comme CHOIX MAÎTRISÉS, pas oublis. Mettre en avant l'API REGISTRE comme vision stratégique. Q : « qu'est-ce qui manque pour la prod ? » → HSM/secrets, audit externe, tests auto + charge, HTTPS durci, cadre RGPD-santé.")

# ============================================================ SLIDE 17 — Conclusion
s = slide(); bg(s, BLEU)
box(s, 0, 0, 13.333, 0.12, fill=VERT, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 7.38, 13.333, 0.12, fill=CYAN, shape=MSO_SHAPE.RECTANGLE)
txt(s, 0.9, 1.5, 11.5, 0.5, [("CONCLUSION", 15, CYAN, True, False)], font=TITLE_FONT)
txt(s, 0.9, 2.2, 11.6, 1.8,
    [("De l'Excel partagé à l'autorité de référence numérique", 32, BLANC, True, False)], font=TITLE_FONT)
txt(s, 0.95, 3.9, 11.4, 1.2,
    [("Une première brique solide du système d'information de l'ONMM : centralisée, sécurisée et traçable — "
      "et une trajectoire claire vers un registre national exposé en API pour tout le secteur de la santé.",
      17, RGBColor(0xC9,0xE6,0xEC), False, False)])
chip(s, 0.95, 5.5, 4.2, 0.8, "Merci — vos questions ?", VERT, size=18)
txt(s, 9.0, 5.55, 3.4, 0.7, [("M. Sidi Mohamed · C23282", 13, RGBColor(0xC9,0xE6,0xEC), False, False)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "CONCLUSION (30s). Finir sur la VISION (ONMM = autorité de référence numérique de la santé), pas sur « voilà j'ai fini ». Remercier le jury, inviter aux questions avec assurance. Phrase de clôture répétée à l'avance. JAMAIS finir par « euh voilà ».")

# ============================================================ ANNEXES
def annex(title, lines):
    s = slide(); bg(s, GRISBG)
    box(s, 0, 0, 13.333, 1.2, fill=ANTHRA, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 0.6, 0.3, 12, 0.7, [[("ANNEXE · ", 16, OR, True, False), (title, 22, BLANC, True, False)]], font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.7, 1.6, 12, 5.2, [[(l, 14, ANTHRA, False, False)] for l in lines])
    return s

annex("Détail cryptographique (sur question vote)",
      ["• Format stocké : Base64(wrappedDEK) + \".\" + Base64(iv + ciphertext)",
       "• AES-256-GCM : tag d'authentification 128 bits, IV 12 octets, DEK 32 octets aléatoire par bulletin",
       "• RSA-2048 OAEP-SHA256 (MGF1-SHA256) : enveloppe (wrap) de la DEK avec la clé publique de l'élection",
       "• Clé privée RSA délivrée par KeyManagementService UNIQUEMENT si statut ≠ VOTE_EN_COURS",
       "• Vote Hash = SHA-256(electionId : voterToken : encryptedChoice : timestamp) — le vote n'apparaît jamais en clair",
       "• Voter Token = HMAC-SHA256(secret serveur, electionId, médecinId) — déterministe (unicité) mais non réversible",
       "• Signature Ed25519 (BouncyCastle) du serveur sur (encryptedChoice || voterToken || timestamp)",
       "• Une seule paire Ed25519 au niveau serveur (origine du backend) ; une paire RSA PAR élection (confidentialité)"])

annex("Machine d'état d'une élection",
      ["BROUILLON  →  VOTE_EN_COURS  →  CLÔTURE  →  DÉPOUILLEMENT  →  RÉSULTATS PUBLIÉS",
       "",
       "• BROUILLON : configuration, candidatures, génération de la paire de clés RSA de l'élection",
       "• VOTE_EN_COURS : bulletins chiffrés + signés acceptés ; clé privée RSA INACCESSIBLE",
       "• CLÔTURE : plus aucun vote accepté",
       "• DÉPOUILLEMENT : déverrouillage de la clé privée, déchiffrement, vérification hash + signature",
       "• RÉSULTATS : agrégation + signature des résultats finaux, publication"])

annex("Planning des 8 sprints (16 semaines)",
      ["Sprint 0 (S1-2)  : Cahier des charges, analyse des besoins, conception UML",
       "Sprint 1 (S3-4)  : Validation administrative, registre officiel, profil médecin, spécialités",
       "Sprint 2 (S5-6)  : Authentification, sécurité, gestion des rôles, module d'adhésion",
       "Sprint 3 (S7-8)  : Pages publiques, annuaire, annonces, diffusion d'information",
       "Sprint 4 (S9-10) : Réclamations et notifications",
       "Sprint 5 (S11-12): Sondages / questionnaires (réponses, UI, statistiques)",
       "Sprint 6 (S13-14): Élections et vote (règles métier, états, sécurité)",
       "Sprint 7 (S15-16): Intégration finale, tests fonctionnels, corrections, déploiement"])

prs.save("/home/user/PFE_2026_DAII_C23282/presentation/Presentation_PFE_ONMM.pptx")
print("OK — slides:", len(prs.slides._sldIdLst))
